"""文檔處理和索引"""
import os
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import uuid

from src.vector_store import VectorStoreManager, HybridSearcher
from src.utils import chunk_text, compute_file_hash

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """文檔處理器"""

    def __init__(
        self,
        vector_store: VectorStoreManager,
        chunk_size: int = 1000,
        chunk_overlap: int = 200
    ):
        """初始化文檔處理器

        Args:
            vector_store: 向量存儲管理器
            chunk_size: 分塊大小
            chunk_overlap: 重疊大小
        """
        self.vector_store = vector_store
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.hybrid_searcher = HybridSearcher(vector_store)

        # 記錄已處理的文件（避免重複處理）
        self.processed_files = {}

    def process_file(self, file_path: str) -> Dict[str, Any]:
        """處理單個文件

        Args:
            file_path: 文件路徑

        Returns:
            處理結果
        """
        try:
            # 檢查文件是否已處理
            file_hash = compute_file_hash(file_path)
            if file_hash in self.processed_files:
                logger.info(f"File already processed: {file_path}")
                return self.processed_files[file_hash]

            # 根據文件類型讀取內容
            file_ext = Path(file_path).suffix.lower()
            content = self._read_file(file_path, file_ext)

            if not content:
                logger.warning(f"No content extracted from {file_path}")
                return {"success": False, "error": "No content extracted"}

            # 分塊
            chunks = chunk_text(content, self.chunk_size, self.chunk_overlap)
            logger.info(f"Split {file_path} into {len(chunks)} chunks")

            # 準備元數據
            filename = Path(file_path).name
            metadatas = [
                {
                    "source": filename,
                    "file_path": file_path,
                    "chunk_index": i,
                    "file_hash": file_hash
                }
                for i in range(len(chunks))
            ]

            # 生成 ID
            doc_ids = [f"{file_hash}_{i}" for i in range(len(chunks))]

            # 添加到向量存儲
            self.vector_store.add_documents(chunks, metadatas, doc_ids)

            # 添加到混合搜索器
            self.hybrid_searcher.add_documents_for_keyword_search(chunks, doc_ids)

            # 記錄處理結果
            result = {
                "success": True,
                "filename": filename,
                "chunks": len(chunks),
                "file_hash": file_hash
            }
            self.processed_files[file_hash] = result

            return result

        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            return {"success": False, "error": str(e)}

    def process_directory(self, directory: str) -> Dict[str, Any]:
        """處理目錄下的所有文件

        Args:
            directory: 目錄路徑

        Returns:
            處理統計
        """
        stats = {
            "total_files": 0,
            "successful": 0,
            "failed": 0,
            "total_chunks": 0,
            "files": []
        }

        supported_extensions = {'.txt', '.md', '.pdf', '.docx', '.pptx'}

        # 遍歷目錄
        for file_path in Path(directory).rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                stats["total_files"] += 1

                result = self.process_file(str(file_path))

                if result.get("success"):
                    stats["successful"] += 1
                    stats["total_chunks"] += result.get("chunks", 0)
                else:
                    stats["failed"] += 1

                stats["files"].append(result)

        logger.info(f"Processed directory: {stats}")
        return stats

    def _read_file(self, file_path: str, file_ext: str) -> str:
        """根據文件類型讀取內容

        Args:
            file_path: 文件路徑
            file_ext: 文件擴展名

        Returns:
            文件內容
        """
        try:
            if file_ext == '.txt' or file_ext == '.md':
                return self._read_text_file(file_path)
            elif file_ext == '.pdf':
                return self._read_pdf(file_path)
            elif file_ext == '.docx':
                return self._read_docx(file_path)
            elif file_ext == '.pptx':
                return self._read_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_ext}")
                return ""
        except Exception as e:
            logger.error(f"Failed to read {file_path}: {e}")
            return ""

    def _read_text_file(self, file_path: str) -> str:
        """讀取文本文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _read_pdf(self, file_path: str) -> str:
        """讀取 PDF 文件"""
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            text = []

            for page in reader.pages:
                text.append(page.extract_text())

            return "\n\n".join(text)

        except ImportError:
            logger.error("PyPDF2 not installed. Install: pip install PyPDF2")
            return ""
        except Exception as e:
            logger.error(f"Failed to read PDF: {e}")
            return ""

    def _read_docx(self, file_path: str) -> str:
        """讀取 Word 文件"""
        try:
            from docx import Document

            doc = Document(file_path)
            text = []

            for paragraph in doc.paragraphs:
                text.append(paragraph.text)

            return "\n\n".join(text)

        except ImportError:
            logger.error("python-docx not installed. Install: pip install python-docx")
            return ""
        except Exception as e:
            logger.error(f"Failed to read DOCX: {e}")
            return ""

    def _read_pptx(self, file_path: str) -> str:
        """讀取 PowerPoint 文件"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

            return "\n\n".join(text)

        except ImportError:
            logger.error("python-pptx not installed. Install: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"Failed to read PPTX: {e}")
            return ""

    def add_text_document(
        self,
        text: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """添加文本文檔

        Args:
            text: 文本內容
            metadata: 元數據

        Returns:
            處理結果
        """
        try:
            # 分塊
            chunks = chunk_text(text, self.chunk_size, self.chunk_overlap)

            # 準備元數據
            doc_id = str(uuid.uuid4())
            metadatas = [
                {
                    **(metadata or {}),
                    "chunk_index": i,
                    "document_id": doc_id
                }
                for i in range(len(chunks))
            ]

            # 生成 ID
            chunk_ids = [f"{doc_id}_{i}" for i in range(len(chunks))]

            # 添加到向量存儲
            self.vector_store.add_documents(chunks, metadatas, chunk_ids)

            # 添加到混合搜索器
            self.hybrid_searcher.add_documents_for_keyword_search(chunks, chunk_ids)

            return {
                "success": True,
                "document_id": doc_id,
                "chunks": len(chunks)
            }

        except Exception as e:
            logger.error(f"Failed to add text document: {e}")
            return {"success": False, "error": str(e)}


def main():
    """命令行工具：初始化文檔索引"""
    import argparse
    from src.utils import load_config, setup_logging

    parser = argparse.ArgumentParser(description="Document indexing tool")
    parser.add_argument(
        "--init",
        action="store_true",
        help="Initialize document index from docs directory"
    )
    parser.add_argument(
        "--docs-dir",
        type=str,
        default="./docs",
        help="Documents directory"
    )
    parser.add_argument(
        "--config",
        type=str,
        default="./config/config.yaml",
        help="Config file path"
    )

    args = parser.parse_args()

    # 加載配置
    config = load_config(args.config)
    setup_logging(config)

    # 初始化組件
    vector_store = VectorStoreManager(
        persist_directory=config.get('vector_store', {}).get('persist_directory', './data/chroma_db'),
        collection_name=config.get('vector_store', {}).get('collection_name', 'documents')
    )

    processor = DocumentProcessor(
        vector_store=vector_store,
        chunk_size=config.get('rag', {}).get('chunk_size', 1000),
        chunk_overlap=config.get('rag', {}).get('chunk_overlap', 200)
    )

    if args.init:
        logger.info(f"Initializing document index from {args.docs_dir}")
        stats = processor.process_directory(args.docs_dir)
        print(f"\n✓ Indexing completed!")
        print(f"  - Total files: {stats['total_files']}")
        print(f"  - Successful: {stats['successful']}")
        print(f"  - Failed: {stats['failed']}")
        print(f"  - Total chunks: {stats['total_chunks']}")
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
